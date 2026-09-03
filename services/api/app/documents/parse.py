"""Parsing uploaded documents into pages with provenance.

Doc 07 M5: *"Chunk with source and page retained"* and *"handle parse failures,
scanned PDFs without OCR, size limits — **visibly, never silently**."*

Provenance is not metadata here, it is the product. Proposal Studio's grounding
rule (doc 01 M8) is that every price cites the document and page it came from,
so a parser that loses the page number quietly breaks the highest-liability
module in the system. Every `Page` therefore carries its number from the moment
it is read.

**A scanned PDF is the failure worth naming.** It parses perfectly and yields no
text. Treated as success it becomes an indexed document with no content — the
user uploaded their price list, saw a tick, and Proposal Studio will later say
it cannot find a price. `ParseOutcome.NO_TEXT_LAYER` exists so that lands as a
visible state at upload instead.
"""

from __future__ import annotations

import io
from dataclasses import dataclass, field
from enum import StrEnum

# One number, defined in `limits.py` with the other two. It lived here at 50 MB
# — `doc/01` M1's figure — after `doc/11` settled on 25, and a second definition
# is how they came to disagree in the first place.
from app.documents.limits import MAX_FILE_BYTES as MAX_FILE_BYTES

# Below this, a "successful" parse is almost certainly an image-only document.
MIN_TEXT_CHARS_PER_PAGE = 20


class DocumentKind(StrEnum):
    PDF = "pdf"
    DOCX = "docx"
    PPTX = "pptx"
    XLSX = "xlsx"
    TXT = "txt"


class ParseOutcome(StrEnum):
    OK = "ok"
    TOO_LARGE = "too_large"
    UNSUPPORTED_TYPE = "unsupported_type"
    CORRUPT = "corrupt"
    ENCRYPTED = "encrypted"
    NO_TEXT_LAYER = "no_text_layer"
    EMPTY = "empty"


# What the user is told. Each names the problem *and* what to do, because a
# failure the user cannot act on is only marginally better than a silent one.
OUTCOME_MESSAGE: dict[ParseOutcome, str] = {
    ParseOutcome.OK: "",
    ParseOutcome.TOO_LARGE: "This file is over 50 MB. Split it and upload the parts.",
    ParseOutcome.UNSUPPORTED_TYPE: "Only PDF, Word, PowerPoint, Excel and text files can be read.",
    ParseOutcome.CORRUPT: "This file could not be opened. It may be damaged.",
    ParseOutcome.ENCRYPTED: "This file is password-protected. Remove the password and re-upload.",
    ParseOutcome.NO_TEXT_LAYER: (
        "This looks like a scan. There is no text to read, so nothing was indexed — "
        "re-save it with text, or use a searchable PDF."
    ),
    ParseOutcome.EMPTY: "This file contains no text.",
}

EXTENSION_KIND: dict[str, DocumentKind] = {
    ".pdf": DocumentKind.PDF,
    ".docx": DocumentKind.DOCX,
    ".pptx": DocumentKind.PPTX,
    ".xlsx": DocumentKind.XLSX,
    ".xlsm": DocumentKind.XLSX,
    ".txt": DocumentKind.TXT,
    ".md": DocumentKind.TXT,
}


@dataclass(frozen=True, slots=True)
class Page:
    """One page, slide, sheet or section — the unit a citation points at."""

    number: int
    text: str
    label: str | None = None
    """A human-meaningful name where the format has one: a sheet name, a slide
    title. "page 4" is far less useful than "Rate Card 2026" in a citation."""


@dataclass(frozen=True, slots=True)
class ParsedDocument:
    kind: DocumentKind | None
    outcome: ParseOutcome
    pages: tuple[Page, ...] = field(default=())
    page_count: int = 0
    char_count: int = 0

    @property
    def succeeded(self) -> bool:
        return self.outcome is ParseOutcome.OK

    @property
    def message(self) -> str:
        return OUTCOME_MESSAGE[self.outcome]


def kind_for_filename(filename: str) -> DocumentKind | None:
    lowered = filename.lower()
    for extension, kind in EXTENSION_KIND.items():
        if lowered.endswith(extension):
            return kind
    return None


def _fail(kind: DocumentKind | None, outcome: ParseOutcome) -> ParsedDocument:
    return ParsedDocument(kind=kind, outcome=outcome)


def parse_document(data: bytes, *, filename: str) -> ParsedDocument:
    """Parse into pages, or return a named failure. Never raises."""
    if len(data) > MAX_FILE_BYTES:
        return _fail(kind_for_filename(filename), ParseOutcome.TOO_LARGE)

    kind = kind_for_filename(filename)
    if kind is None:
        return _fail(None, ParseOutcome.UNSUPPORTED_TYPE)

    if not data:
        return _fail(kind, ParseOutcome.EMPTY)

    try:
        pages = _PARSERS[kind](data)
    except _EncryptedError:
        return _fail(kind, ParseOutcome.ENCRYPTED)
    except Exception:
        # Deliberately broad. A malformed file is a user-visible outcome, not a
        # 500 — and the parsers raise a wide variety of library-specific errors.
        return _fail(kind, ParseOutcome.CORRUPT)

    total_chars = sum(len(p.text.strip()) for p in pages)

    if not pages or total_chars == 0:
        # PDFs are the common case here; a Word file with no text is simply
        # empty, which is a different thing to say.
        outcome = ParseOutcome.NO_TEXT_LAYER if kind is DocumentKind.PDF else ParseOutcome.EMPTY
        return ParsedDocument(kind=kind, outcome=outcome, page_count=len(pages))

    if kind is DocumentKind.PDF and total_chars < MIN_TEXT_CHARS_PER_PAGE * len(pages):
        # Scanned PDFs frequently yield a handful of stray characters rather
        # than nothing at all, which would otherwise pass as a real parse.
        return ParsedDocument(
            kind=kind,
            outcome=ParseOutcome.NO_TEXT_LAYER,
            page_count=len(pages),
            char_count=total_chars,
        )

    return ParsedDocument(
        kind=kind,
        outcome=ParseOutcome.OK,
        pages=tuple(pages),
        page_count=len(pages),
        char_count=total_chars,
    )


class _EncryptedError(Exception):
    pass


def _parse_pdf(data: bytes) -> list[Page]:
    from pypdf import PdfReader
    from pypdf.errors import FileNotDecryptedError

    reader = PdfReader(io.BytesIO(data))
    if reader.is_encrypted:
        try:
            # An empty-password PDF is encrypted but readable.
            if reader.decrypt("") == 0:
                raise _EncryptedError
        except FileNotDecryptedError as exc:
            raise _EncryptedError from exc

    return [
        Page(number=index + 1, text=(page.extract_text() or ""))
        for index, page in enumerate(reader.pages)
    ]


def _parse_docx(data: bytes) -> list[Page]:
    from docx import Document

    document = Document(io.BytesIO(data))
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]

    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))

    # Word has no page concept before rendering, so the document is one logical
    # page. Claiming page numbers we cannot know would make citations wrong in
    # a way nobody could detect.
    return [Page(number=1, text="\n".join(paragraphs), label="document")]


def _parse_pptx(data: bytes) -> list[Page]:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    pages: list[Page] = []

    for index, slide in enumerate(presentation.slides):
        parts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        title = None
        if slide.shapes.title is not None and slide.shapes.title.text.strip():
            title = slide.shapes.title.text.strip()[:120]
        pages.append(Page(number=index + 1, text="\n".join(parts), label=title))

    return pages


def _parse_xlsx(data: bytes) -> list[Page]:
    from openpyxl import load_workbook

    # `data_only` reads computed values rather than formulas: a price list
    # citation should quote the price, not `=B2*1.05`.
    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    pages: list[Page] = []

    for index, sheet in enumerate(workbook.worksheets):
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append(" | ".join(cells))
        pages.append(Page(number=index + 1, text="\n".join(rows), label=sheet.title))

    workbook.close()
    return pages


def _parse_txt(data: bytes) -> list[Page]:
    return [Page(number=1, text=data.decode("utf-8", errors="replace"), label="document")]


_PARSERS = {
    DocumentKind.PDF: _parse_pdf,
    DocumentKind.DOCX: _parse_docx,
    DocumentKind.PPTX: _parse_pptx,
    DocumentKind.XLSX: _parse_xlsx,
    DocumentKind.TXT: _parse_txt,
}
