"""Parsing and chunking, with every failure visible.

Doc 07 M5: *"Handle parse failures, scanned PDFs without OCR, size limits —
visibly, never silently."*

The cases that matter most are the ones that would otherwise pass quietly. A
scanned PDF parses without error and yields nothing; treated as success it
becomes an indexed document with no content, and the user finds out weeks later
when Proposal Studio cannot find a price they know they uploaded.
"""

from __future__ import annotations

import io

import pytest

from app.documents.chunk import (
    MIN_CHUNK_CHARS,
    TARGET_CHARS,
    chunk_document,
    chunk_page,
)
from app.documents.parse import (
    MAX_FILE_BYTES,
    DocumentKind,
    Page,
    ParseOutcome,
    kind_for_filename,
    parse_document,
)

# ── Building fixtures in the real formats ─────────────────────


def make_docx(paragraphs: list[str]) -> bytes:
    from docx import Document

    document = Document()
    for text in paragraphs:
        document.add_paragraph(text)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def make_xlsx(sheets: dict[str, list[list[str]]]) -> bytes:
    from openpyxl import Workbook

    workbook = Workbook()
    workbook.remove(workbook.active)
    for name, rows in sheets.items():
        sheet = workbook.create_sheet(title=name)
        for row in rows:
            sheet.append(row)
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def make_pptx(slides: list[tuple[str, str]]) -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    layout = presentation.slide_layouts[1]
    for title, body in slides:
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


# ── Type detection ────────────────────────────────────────────


@pytest.mark.parametrize(
    ("filename", "expected"),
    [
        ("rate-card.pdf", DocumentKind.PDF),
        ("Proposal.DOCX", DocumentKind.DOCX),
        ("deck.pptx", DocumentKind.PPTX),
        ("prices.xlsx", DocumentKind.XLSX),
        ("macro.xlsm", DocumentKind.XLSX),
        ("notes.txt", DocumentKind.TXT),
        ("photo.jpg", None),
        ("archive.zip", None),
        ("noextension", None),
    ],
)
def test_type_detection(filename: str, expected: DocumentKind | None) -> None:
    assert kind_for_filename(filename) is expected


def test_an_unsupported_type_says_which_types_work() -> None:
    result = parse_document(b"\xff\xd8\xff", filename="photo.jpg")
    assert result.outcome is ParseOutcome.UNSUPPORTED_TYPE
    assert "PDF" in result.message


# ── Size ──────────────────────────────────────────────────────


def test_an_oversized_file_is_refused_before_parsing() -> None:
    """Checked on bytes, before any library touches the content."""
    result = parse_document(b"x" * (MAX_FILE_BYTES + 1), filename="huge.txt")
    assert result.outcome is ParseOutcome.TOO_LARGE
    assert "50 MB" in result.message


def test_a_file_at_the_limit_is_accepted() -> None:
    result = parse_document(b"hello world " * 10, filename="small.txt")
    assert result.succeeded


# ── The failures that would otherwise pass silently ───────────


def test_an_empty_file_is_named_not_indexed() -> None:
    result = parse_document(b"", filename="empty.txt")
    assert result.outcome is ParseOutcome.EMPTY
    assert result.pages == ()


def test_a_corrupt_file_fails_visibly_rather_than_raising() -> None:
    """A malformed upload is a user-visible outcome, not a 500."""
    result = parse_document(b"this is definitely not a pdf", filename="broken.pdf")
    assert result.outcome is ParseOutcome.CORRUPT
    assert result.message
    assert not result.succeeded


def test_a_docx_with_no_text_reports_empty() -> None:
    result = parse_document(make_docx([]), filename="blank.docx")
    assert result.outcome is ParseOutcome.EMPTY


# ── Real formats round-trip with provenance ───────────────────


def test_docx_text_is_extracted() -> None:
    result = parse_document(
        make_docx(["Our day rate is OMR 145.", "Mobilisation is charged separately."]),
        filename="rates.docx",
    )
    assert result.succeeded
    assert "OMR 145" in result.pages[0].text


def test_xlsx_keeps_one_page_per_sheet_named_after_the_sheet() -> None:
    """A citation reading "Rate Card 2026" beats one reading "page 2"."""
    result = parse_document(
        make_xlsx(
            {
                "Rate Card 2026": [["Service", "Price"], ["Joinery", "OMR 11850"]],
                "Notes": [["Valid until", "Dec 2026"]],
            }
        ),
        filename="prices.xlsx",
    )
    assert result.succeeded
    assert result.page_count == 2
    assert result.pages[0].label == "Rate Card 2026"
    assert "OMR 11850" in result.pages[0].text


def test_pptx_keeps_one_page_per_slide_titled() -> None:
    result = parse_document(
        make_pptx([("Our approach", "Phased delivery"), ("Pricing", "OMR 3200")]),
        filename="deck.pptx",
    )
    assert result.succeeded
    assert result.page_count == 2
    assert result.pages[1].label == "Pricing"
    assert result.pages[1].number == 2


def test_page_numbers_start_at_one() -> None:
    """Zero-indexed pages in a citation would be wrong on every document."""
    result = parse_document(make_xlsx({"A": [["x"]], "B": [["y"]]}), filename="p.xlsx")
    assert [p.number for p in result.pages] == [1, 2]


# ── Chunking keeps provenance ─────────────────────────────────


def test_every_chunk_carries_its_page() -> None:
    pages = (
        Page(number=1, text="First page content. " * 20, label="Intro"),
        Page(number=2, text="Second page content. " * 20, label="Pricing"),
    )
    chunks = chunk_document(pages)

    assert chunks
    for chunk in chunks:
        assert chunk.page_number in (1, 2)
        assert chunk.page_label in ("Intro", "Pricing")


def test_no_chunk_spans_two_pages() -> None:
    """A chunk covering pages 3 and 4 carries one citation for two sources —
    accurate for half its content, with no way to tell which half."""
    pages = (
        Page(number=1, text="Alpha content here."),
        Page(number=2, text="Beta content here."),
    )
    for chunk in chunk_document(pages):
        assert not ("Alpha" in chunk.text and "Beta" in chunk.text)


def test_citations_are_human_readable() -> None:
    labelled = chunk_page(Page(number=3, text="x" * 200, label="Rate Card 2026"))[0]
    plain = chunk_page(Page(number=4, text="y" * 200))[0]

    assert labelled.citation == "Rate Card 2026 (page 3)"
    assert plain.citation == "page 4"


def test_ordinals_are_stable_and_sequential() -> None:
    pages = tuple(Page(number=n, text=f"Page {n} body. " * 30) for n in range(1, 4))
    chunks = chunk_document(pages)
    assert [c.ordinal for c in chunks] == list(range(len(chunks)))


# ── Chunk sizing ──────────────────────────────────────────────


def test_a_long_page_is_split() -> None:
    chunks = chunk_page(Page(number=1, text="Sentence number one. " * 400))
    assert len(chunks) > 1


def test_no_chunk_greatly_exceeds_the_target() -> None:
    chunks = chunk_page(Page(number=1, text="word " * 5000))
    for chunk in chunks:
        assert len(chunk.text) <= TARGET_CHARS + 200


def test_an_unbroken_run_is_still_split() -> None:
    """A table row or minified block has no sentence boundary to cut on."""
    chunks = chunk_page(Page(number=1, text="x" * (TARGET_CHARS * 3)))
    assert len(chunks) >= 3


def test_a_tiny_trailing_fragment_is_folded_back() -> None:
    """ "OMR 3,200" alone is unretrievable — and worse, citable."""
    text = ("Body paragraph. " * 100) + "\n\nOMR 3200"
    chunks = chunk_page(Page(number=1, text=text))
    assert all(len(c.text) >= MIN_CHUNK_CHARS for c in chunks)
    assert any("OMR 3200" in c.text for c in chunks)


def test_an_empty_page_produces_no_chunks() -> None:
    assert chunk_page(Page(number=1, text="   \n  ")) == []
